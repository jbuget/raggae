import logging
from collections.abc import Callable
from io import BytesIO
from typing import TYPE_CHECKING

from raggae.domain.exceptions.document_exceptions import DocumentExtractionError

if TYPE_CHECKING:
    from raggae.infrastructure.services.pdf_table_extractor import PdfTableExtractor, TableData

logger = logging.getLogger(__name__)


class MultiFormatDocumentTextExtractor:
    """Extract text from txt/md/pdf/docx/pptx files with basic normalization."""

    def __init__(self, pdf_table_extractor: "PdfTableExtractor | None" = None) -> None:
        from raggae.infrastructure.services.pdf_table_extractor import PdfTableExtractor as _PTE

        self._pdf_table_extractor = pdf_table_extractor or _PTE()

    async def extract_text(self, file_name: str, content: bytes, content_type: str) -> str:
        _ = content_type
        extension = file_name.rsplit(".", maxsplit=1)[-1].lower() if "." in file_name else ""

        if extension in {"txt", "md"}:
            text = self._decode_text(content)
        elif extension == "pdf":
            text = self._extract_pdf(content)
        elif extension == "docx":
            text = self._extract_docx(content)
        elif extension == "doc":
            raise DocumentExtractionError("DOC extraction is not supported in sync mode yet. Use DOCX.")
        elif extension == "pptx":
            text = self._extract_pptx(content)
        elif extension == "ppt":
            raise DocumentExtractionError(
                "PPT (legacy binary format) is not supported. Convert to PPTX first."
            )
        else:
            raise DocumentExtractionError(f"Unsupported extension for extraction: {extension}")

        normalized = self._normalize_text(text)
        if not normalized:
            raise DocumentExtractionError("No extractable text found in document")
        return normalized

    def _decode_text(self, content: bytes) -> str:
        try:
            return content.decode("utf-8")
        except UnicodeDecodeError:
            return content.decode("utf-8", errors="ignore")

    def _extract_pdf(self, content: bytes) -> str:
        try:
            import pdfplumber
        except ModuleNotFoundError as exc:  # pragma: no cover
            raise DocumentExtractionError("pdfplumber is required for PDF extraction") from exc

        tables: list[TableData] = []
        try:
            tables = self._pdf_table_extractor.extract_tables(content)
        except Exception:
            logger.warning("pdf_table_extraction_failed", exc_info=True)

        bboxes_by_page: dict[int, list[tuple[float, float, float, float]]] = {}
        for t in tables:
            if t.bbox is not None:
                bboxes_by_page.setdefault(t.page_number, []).append(t.bbox)

        tables_by_page: dict[int, list[TableData]] = {}
        for t in tables:
            tables_by_page.setdefault(t.page_number, []).append(t)

        try:
            parts: list[str] = []
            with pdfplumber.open(BytesIO(content)) as pdf:
                for page_num, page in enumerate(pdf.pages, start=1):
                    page_bboxes = bboxes_by_page.get(page_num, [])
                    if page_bboxes:
                        filtered = page.filter(self._make_bbox_filter(page_bboxes))
                        text = filtered.extract_text() or ""
                    else:
                        text = page.extract_text() or ""
                    for table in tables_by_page.get(page_num, []):
                        parts.append(table.to_chunk())
                    parts.append(f"[[PAGE:{page_num}]]\n{text}")

            return "\n".join(parts)
        except Exception as exc:
            raise DocumentExtractionError(f"Failed to extract PDF text: {exc}") from exc

    @staticmethod
    def _make_bbox_filter(
        bboxes: list[tuple[float, float, float, float]],
    ) -> Callable[[dict[str, object]], bool]:

        def keep(obj: dict[str, object]) -> bool:
            for b in bboxes:
                if (
                    float(obj.get("x0", 0) or 0) >= b[0]  # type: ignore[arg-type]
                    and float(obj.get("x1", 0) or 0) <= b[2]  # type: ignore[arg-type]
                    and float(obj.get("top", 0) or 0) >= b[1]  # type: ignore[arg-type]
                    and float(obj.get("bottom", 0) or 0) <= b[3]  # type: ignore[arg-type]
                ):
                    return False
            return True

        return keep

    def _extract_docx(self, content: bytes) -> str:
        try:
            from docx import Document as DocxDocument
        except ModuleNotFoundError as exc:  # pragma: no cover - dependency management
            raise DocumentExtractionError("python-docx is required for DOCX extraction") from exc

        try:
            document = DocxDocument(BytesIO(content))
            parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
            for table in document.tables:
                for row in table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_cells:
                        parts.append(" | ".join(row_cells))
            return "\n".join(parts)
        except Exception as exc:  # pragma: no cover - file dependent
            raise DocumentExtractionError(f"Failed to extract DOCX text: {exc}") from exc

    def _extract_pptx(self, content: bytes) -> str:
        try:
            from pptx import Presentation
        except ModuleNotFoundError as exc:  # pragma: no cover - dependency management
            raise DocumentExtractionError("python-pptx is required for PPTX extraction") from exc

        try:
            prs = Presentation(BytesIO(content))
            slide_parts: list[str] = []
            for index, slide in enumerate(prs.slides):
                parts: list[str] = [f"[SLIDE:{index + 1}]"]
                title_shape = slide.shapes.title
                title_text = (
                    title_shape.text_frame.text.strip() if title_shape and title_shape.has_text_frame else ""
                )
                if title_text:
                    parts.append(f"# {title_text}")

                sorted_shapes = sorted(
                    slide.shapes,
                    key=lambda s: (getattr(s, "top", 0) or 0, getattr(s, "left", 0) or 0),
                )
                for shape in sorted_shapes:
                    if shape is title_shape:
                        continue
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            parts.append(text)
                    elif shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if cells:
                                parts.append("| " + " | ".join(cells) + " |")

                try:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                except Exception as exc:
                    logger.debug("Impossible de lire les notes de la diapositive %d : %s", index + 1, exc)
                    notes_text = ""
                if notes_text:
                    parts.append(f"[NOTES]\n{notes_text}")

                slide_parts.append("\n".join(parts))
            return "\n".join(slide_parts)
        except Exception as exc:  # pragma: no cover - file dependent
            raise DocumentExtractionError(f"Failed to extract PPTX text: {exc}") from exc

    def _normalize_text(self, text: str) -> str:
        return "\n".join(line.rstrip() for line in text.replace("\r\n", "\n").split("\n")).strip()
